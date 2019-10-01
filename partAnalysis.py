import os
import sys
from classes import *
from plotFuncs import *
import numpy as np
from pptx import Presentation
from pptx.util import Inches
from openpyxl.styles import Alignment
from openpyxl.styles.borders import Border, Side
from openpyxl.styles import PatternFill

start_month = int(sys.argv[1])
start_year = int(sys.argv[2])
end_month = int(sys.argv[3])
end_year = int(sys.argv[4])

if (start_year > end_year) or \
   (start_year == end_year and start_month > end_month):
    print('Podany zakres dat jest błędny')
    sys.exit()
else:
    # generating list of sheetnames for a given range
    list_of_sheetnames = []
    n_of_months = 12*(end_year-start_year) + end_month-start_month + 1
    month = start_month
    year = start_year
    for i in range(n_of_months):
        name_of_sheet = ('0'+str(month) if month < 9 else str(month)) \
                        + '.' + str(year)
        list_of_sheetnames.append(name_of_sheet)

        month += 1
        if month > 12:
            month -= 12
            year += 1

# 1. Loading the file with data for total period
file_path = os.getcwd() + '/data/total.xlsx'
part_label = str(start_month) + '.20' + str(start_year) + '-' \
             + str(end_month) + '.20' + str(end_year)

myWorkbook = MyWorkbook(file_path, list_of_sheetnames)

start_label = [int(list_of_sheetnames[0][:2]), int(list_of_sheetnames[0][-2:])]


# 2. Preparing data from the parsed sheets for visualization
# -- Part as a whole
# a) Barplot for all categories
index_order = np.flip(np.argsort(myWorkbook.cats_sums_list), axis=0)
values_desc = [myWorkbook.cats_sums_list[i] for i in index_order
               if myWorkbook.cats_sums_list[i] > 0]
labels_desc = [myWorkbook.cats_names[i] for i in index_order
               if myWorkbook.cats_sums_list[i] > 0]

# b) Piechart of spendings for the main categories
index_order = np.flip(np.argsort(myWorkbook.cats_sums_list), axis=0)
_top_indices = index_order[0:5]
_low_indices = index_order[5:len(index_order)]

top_labels = [myWorkbook.cats_names[i] for i in _top_indices]
top_values = [myWorkbook.cats_sums_list[i] for i in _top_indices]
others_values = [myWorkbook.cats_sums_list[i] for i in _low_indices]

top_plus_others_values = top_values + [sum(others_values)]
top_plus_others_labels_0 = top_labels + ['Pozostałe']
top_plus_others_labels = [top_plus_others_labels_0[i] + ' - '
                          + str(round(top_plus_others_values[i], 2)) + ' zł'
                          for i in range(len(top_plus_others_values))]

# c) Piechart of the metacategories
metacats_values = [myWorkbook.sum_basic, myWorkbook.sum_addit,
                   myWorkbook.sum_giftdon]
metacats_labels = ['Podstawowe - ' + str(round(myWorkbook.sum_basic, 2))
                   + 'zł', 'Dodatkowe - ' + str(round(myWorkbook.sum_addit, 2))
                   + 'zł', 'Prezenty i donacje - '
                   + str(round(myWorkbook.sum_giftdon, 2)) + 'zł']

# d) Piechart of incomes
_values_list_inc = list(myWorkbook.incomes_dict.values())
_labels_list_inc = list(myWorkbook.incomes_dict.keys())
incomes_values = [inc for inc in _values_list_inc if inc > 0]
incomes_labels = [inc+' - '+str(_values_list_inc[i])+'zł' for i, inc in
                  enumerate(_labels_list_inc) if _values_list_inc[i] > 0]

# e) Piechart of earnings
_values_list_ear = list(myWorkbook.earnings_dict.values())
_labels_list_ear = list(myWorkbook.earnings_dict.keys())
earnings_values = [ear for ear in _values_list_ear if ear > 0]
earnings_labels = [ear + ' - ' + str(_values_list_ear[i])+'zł' for i, ear in
                   enumerate(_labels_list_ear) if _values_list_ear[i] > 0]

# -- Averaged month --
# f) Piechart of spendings for the main categories
top_plus_others_values_avg = [i/n_of_months for i in top_plus_others_values]
top_plus_others_labels_avg = [top_plus_others_labels_0[i] + ' - '
                              + str(round(top_plus_others_values[i]
                                          / n_of_months, 2)) + ' zł'
                              for i in range(len(top_plus_others_values))]

# g) Piechart of the metacategories
metacats_values_avg = [i/n_of_months for i in metacats_values]
metacats_labels_avg = ['Podstawowe - '
                       + str(round(myWorkbook.sum_basic/n_of_months, 2))
                       + 'zł', 'Dodatkowe - '
                       + str(round(myWorkbook.sum_addit/n_of_months, 2))
                       + 'zł', 'Prezenty i donacje - '
                       + str(round(myWorkbook.sum_giftdon/n_of_months, 2))
                       + 'zł']

# h) Piechart of incomes
incomes_values_avg = [i/n_of_months for i in incomes_values]
incomes_labels_avg = [inc + ' - ' + str(round(_values_list_inc[i]
                      / n_of_months, 2)) + 'zł'
                      for i, inc in enumerate(_labels_list_inc)
                      if _values_list_inc[i] > 0]

# i) Piechart of earnings
earnings_values_avg = [i/n_of_months for i in earnings_values]
earnings_labels_avg = [ear + ' - '
                       + str(round(_values_list_ear[i]/n_of_months, 2))
                       + 'zł'
                       for i, ear in enumerate(_labels_list_ear)
                       if _values_list_ear[i] > 0]

# -- Part as a sequence of months --
# j) Stackplot of cummulated spendings for the top categories
top_spends_seqs = []
# top categories
for c, cat in enumerate(top_labels):
    top_spends_seqs.append([])
    for month in myWorkbook.sheets_list:
        top_spends_seqs[c].append(month.cats_sums[cat])
# others
top_spends_seqs.append([])
for m, month in enumerate(myWorkbook.sheets_list):
    top_spends_seqs[-1].append(np.sum(month.cats_sums_list)
                               - sum(cat[m] for cat in top_spends_seqs[:-1]
                               if cat[m]))
# cummulated sums of the lists
top_spends_seqs_cumsum = []
for seq in top_spends_seqs:
    top_spends_seqs_cumsum.append(np.cumsum(seq))

# k) Lineplot of cummulated spendings, incomes and savings
line_incomes = []
line_spendings = []
line_balance = []
line_savings = []

for month in myWorkbook.sheets_list:
    line_incomes.append(month.incomes)
    line_spendings.append(sum(month.cats_sums_list))
    line_balance.append(month.balance[0])
    line_savings.append(month.savings)

line_incomes = np.cumsum(line_incomes)
line_spendings = np.cumsum(line_spendings)
line_balance = np.cumsum(line_balance)
line_savings = np.cumsum(line_savings)

# l) Lineplot of spendings and incomes in subsequent months
spendings_list = []
incomes_list = []
for month in myWorkbook.sheets_list:
    spendings_list.append(month.sum_total)
    incomes_list.append(month.incomes)

# m) Lineplot of average spendings for subsequent categories so far
current_means_seqs = []
for c, cat in enumerate(top_spends_seqs):
    current_means_seqs.append([])
    for m, month in enumerate(cat):
        current_means_seqs[-1].append(np.mean(top_spends_seqs[c][:m+1]))

# n) Lineplot of main sources
#  choosing of the relevant sources
incomes_main = []
for inc in myWorkbook.incomes_dict:
    if myWorkbook.incomes_dict[inc] > 0.02*myWorkbook.incomes:
        incomes_main.append(inc)

# preprocessing
incomes_seqs = [[] for i in range(len(incomes_main))]
for p, inc in enumerate(incomes_main):
    for m, month in enumerate(myWorkbook.sheets_list):
        try:
            incomes_seqs[p].append(month.incomes_dict[inc])
        except:
            incomes_seqs[p].append(0)


# 3. Visualization and saving the plots
results_dir = os.getcwd() + '/results/' + part_label + ' - wyniki'
if not os.path.exists(results_dir):
    os.mkdir(results_dir)
    os.mkdir(results_dir + '/plots')

# a) Barplot for all categories
values = values_desc
labels = labels_desc
title = part_label + ' - Kwoty wydane w ciągu całego okresu \nna ' \
                      'kolejne kategorie'
fig_name = results_dir + '/plots/plot1.png'

fig = plotBar(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# b) Piechart of spendings for the main categories
values = top_plus_others_values
labels = top_plus_others_labels
title = part_label \
        + ' - Struktura wydatków z podziałem \nna kategorie\n\n' \
        'Suma wydatków: ' + str(round(myWorkbook.sum_total, 2)) + 'zł\n'
fig_name = results_dir + '/plots/plot2.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# c) Piechart of the metacategories
values = metacats_values
labels = metacats_labels
title = part_label + ' - Podział wydatków na:\n' \
        + 'Podstawowe, Dodatkowe i Prezenty/Donacje\n\nSuma wydatków: ' \
        + str(round(myWorkbook.sum_total, 2)) + 'zł\n'
fig_name = results_dir + '/plots/plot3.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# d) Piechart of incomes
values = incomes_values
labels = incomes_labels
title = part_label + ' - Podział przychodów na \nposzczególne źródła\n\n' \
        'Suma przychodów: ' + str(round(myWorkbook.incomes, 2)) + 'zł\n' \
        'Nadwyżka przychodów: ' + str(round(myWorkbook.balance[0], 2)) \
        + 'zł  (' + str(round(100*myWorkbook.balance[0]
                               / myWorkbook.incomes, 2)) \
        + '%)'

fig_name = results_dir + '/plots/plot4.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# e) Piechart of earnings
values = earnings_values
labels = earnings_labels
title = part_label + ' - Podział zarobków na \nposzczególne źrodła\n\n' \
        + 'Suma zarobków: ' + str(myWorkbook.earnings) + 'zł\n' \
        + 'Nadwyżka zarobków: ' + str(round(myWorkbook.balance[1], 2))\
        + 'zł  (' + str(round(100*myWorkbook.balance[1]
                               / myWorkbook.earnings, 2)) \
        + '%)'
fig_name = results_dir + '/plots/plot5.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# f) Piechart of spendings for the main categories
values = top_plus_others_values_avg
labels = top_plus_others_labels_avg
title = part_label \
        + ' - Struktura wydatków w uśrednionym \nmiesiącu z podziałem na ' \
        'kategorie\n\n'+'Suma wydatków: ' \
        + str(round(myWorkbook.sum_total/n_of_months, 2)) + 'zł\n'
fig_name = results_dir + '/plots/plot6.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# g) Piechart of metacategories
values = metacats_values_avg
labels = metacats_labels_avg
title = part_label + ' - Podział wydatków na: ' \
        + 'Podstawowe, \nDodatkowe i Prezenty/Donacje ' \
        'w uśrednionym miesiącu\n\n' \
        'Suma wydatków: ' + str(round(myWorkbook.sum_total/n_of_months, 2)) \
        + 'zł\n'
fig_name = results_dir + '/plots/plot7.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# h) Piechart of incomes
values = incomes_values_avg
labels = incomes_labels_avg
title = part_label + ' - Podział przychodów na poszczególne \nźródła ' \
        'w uśrednionym miesiącu\n\n' + 'Suma przychodów: ' \
        + str(round(myWorkbook.incomes/n_of_months, 2)) + 'zł\n' \
        + 'Nadwyżka przychodów: ' \
        + str(round(myWorkbook.balance[0]/n_of_months, 2)) \
        + 'zł  (' + str(round(100*myWorkbook.balance[0] / n_of_months
                               / (myWorkbook.incomes/n_of_months), 2)) + '%)'

fig_name = results_dir + '/plots/plot8.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# i) Piechart of earnings
values = earnings_values_avg
labels = earnings_labels_avg
title = part_label + ' - Podział zarobków na poszczególne \nźródła ' \
        'w uśrednionym miesiącu\n\nSuma zarobków: ' \
        + str(round(myWorkbook.earnings/n_of_months, 2)) + 'zł\n' \
        'Nadwyżka zarobków: ' \
        + str(round(myWorkbook.balance[1]/n_of_months, 2)) \
        + 'zł  (' + str(round(100*myWorkbook.balance[1]/n_of_months
                               / (myWorkbook.earnings/n_of_months), 2)) \
        + '%)'

fig_name = results_dir + '/plots/plot9.png'

fig = plotPie(values, labels, title)
plt.savefig(figure=fig, fname=fig_name)

# j) Stackplot of cummulated spendings for the top categories
values = top_spends_seqs_cumsum
labels = top_labels + ['Pozostałe']
title = part_label + ' - Skumulowane wartości wydatków na\n poszczególne ' \
                      'kategorie'
fig_name = results_dir + '/plots/plot10.png'

fig = plotStack(values, labels, title, start_label)
plt.savefig(figure=fig, fname=fig_name)

# k) Lineplot of cummulated spendings, incomes and savings
values = [line_incomes, line_spendings, line_balance, line_savings]
labels = ['Przychody', 'Wydatki', 'Nadwyżka\nprzychodów',
          'Oszczędności\ndługoterminowe']
title = part_label + ' - Skumulowane wartości przychodów, wydatków \n' \
        'i oszczędności'
fig_name = results_dir + '/plots/plot11.png'

fig = plotLine(values, labels, title, start_label)
plt.savefig(figure=fig, fname=fig_name)

# l) Lineplot of spendings and incomes in subsequent months
values = [incomes_list, spendings_list]
labels = ['Przychody', 'Wydatki']
title = part_label + ' - Przychody i wydatki w kolejnych miesiącach'
fig_name = results_dir + '/plots/plot12.png'

fig = plotLine(values, labels, title, start_label)
plt.savefig(figure=fig, fname=fig_name)

# m) Lineplot of average spendings for subsequent categories so far
values = current_means_seqs
labels = top_labels + ['Pozostałe']
title = part_label \
        + ' - Dotychczasowe średnie miesięczne wydatki na \nposzczególne ' \
        'kategorie'
fig_name = results_dir + '/plots/plot13.png'

fig = plotLine(values, labels, title, start_label)
plt.savefig(figure=fig, fname=fig_name)

# n) Lineplot of the sources
values = incomes_seqs
labels = incomes_main
title = part_label + ' - Kwoty przychodów z najważniejszych \n źrodeł'
fig_name = results_dir + '/plots/plot14.png'

fig = plotLine(values, labels, title, start_label)
plt.savefig(figure=fig, fname=fig_name)

plt.close('all')

# 4. Pptx presentation
prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
blank_slide_layout = prs.slide_layouts[6]
slides = list()

slides.append(prs.slides.add_slide(title_slide_layout))
title = slides[-1].shapes.title
title.text = part_label + ' - raport finansowy'

# Part as a whole
slides.append(prs.slides.add_slide(title_slide_layout))
title = slides[-1].placeholders[1]
title.text = '1. Cały przedział'

for i in range(5):
    slides.append(prs.slides.add_slide(blank_slide_layout))
    left = Inches(1.5)
    top = Inches(0.0)
    height = width = Inches(7.5)

    pic_path = 'results/' + part_label + ' - wyniki/plots/plot' \
               + str(i+1) + '.png'
    pic = slides[-1].shapes.add_picture(pic_path, left, top, height, width)

# Averaged month
slides.append(prs.slides.add_slide(title_slide_layout))
title = slides[-1].placeholders[1]
title.text = '2. Uśredniony miesiąc'

for i in range(4):
    slides.append(prs.slides.add_slide(blank_slide_layout))
    left = Inches(1.5)
    top = Inches(0.0)
    height = width = Inches(7.5)

    pic_path = 'results/' + part_label + ' - wyniki/plots/plot' \
               + str(i+1+5) + '.png'
    pic = slides[-1].shapes.add_picture(pic_path, left, top, height, width)

# Part as a sequence of months
slides.append(prs.slides.add_slide(title_slide_layout))
title = slides[-1].placeholders[1]
title.text = '3. Sekwencja miesięcy'


for i in range(5):
    slides.append(prs.slides.add_slide(blank_slide_layout))
    left = Inches(0.0)
    top = Inches(0.1)
    height = Inches(10.5)
    width = Inches(7.0)

    pic_path = 'results/' + part_label + ' - wyniki/plots/plot' \
               + str(i+1+9) + '.png'
    pic = slides[-1].shapes.add_picture(pic_path, left, top, height, width)


prs.save(results_dir + '/' + part_label + ' - raport finansowy.pptx')

# 5. Exporting spendings to a table in seprate excel workbook
wb_to_export = openpyxl.Workbook()


def export_to_excel(cat_name, num_of_ws):
    values = myWorkbook.spends_values_yr[cat_name]
    items = myWorkbook.spends_items_yr[cat_name]
    monthlabels = myWorkbook.spends_monthlabel_yr[cat_name]

    mdict = {1: 'Styczeń', 2: 'Luty', 3: 'Marzec', 4: 'Kwiecień', 5: 'Maj',
             6: 'Czerwiec', 7: 'Lipiec', 8: 'Sierpień', 9: 'Wrzesień',
             10: 'Październik', 11: 'Listopad', 12: 'Grudzień'}

    values_dict = dict()
    items_dict = dict()

    for month_num in np.unique(monthlabels):
        _one_month_list = [values[i] for i, j in enumerate(monthlabels)
                           if j == month_num]
        values_dict[month_num] = _one_month_list

        _one_month_list = [items[i] for i, j in enumerate(monthlabels)
                           if j == month_num]
        items_dict[month_num] = _one_month_list

    ws = wb_to_export.create_sheet(cat_name, num_of_ws)

    # 'try' becase doesn't work when no spendings
    try:
        ws.column_dimensions['A'].width = max([len(i) for i in items]) + 2
    except:
        pass

    row = 0
    month = start_label[0]
    year = start_label[1]
    for number in range(1, n_of_months+1):     # list(values_dict.keys()):
        if number in monthlabels:
            row += 1
            ws.cell(row, 1).value = mdict[month] + '  ' + str(month) \
                                                 + '.20' + str(year)
            ws.merge_cells(start_row=row, start_column=1, end_row=row,
                           end_column=2)
            ws.cell(row, 1).alignment = Alignment(horizontal='center')
            ws.cell(row, 1).fill = PatternFill(fgColor='93e1e6',
                                               fill_type='solid')

            row += 1
            ws.cell(row, 1).value = 'Co?'
            ws.cell(row, 2).value = 'Kwota [zł]'
            ws.cell(row, 1).alignment = Alignment(horizontal='center')
            ws.cell(row, 2).alignment = Alignment(horizontal='center')
            ws.cell(row, 1).fill = PatternFill(fgColor='daf6f7',
                                               fill_type='solid')
            ws.cell(row, 2).fill = PatternFill(fgColor='daf6f7',
                                               fill_type='solid')

            for i, val in enumerate(values_dict[number]):
                row += 1
                ws.cell(row, 1).value = items_dict[number][i]
                ws.cell(row, 2).value = val

        month += 1
        if month > 12:
            month -= 12
            year += 1

    thin_border = Border(left=Side(style='thin'),
                         right=Side(style='thin'),
                         top=Side(style='thin'),
                         bottom=Side(style='thin'))
    for r in range(row):
        ws.cell(row=r + 1, column=1).border = thin_border
        ws.cell(row=r + 1, column=2).border = thin_border

    return wb_to_export


wb_to_export.remove_sheet(wb_to_export.active)
wb_to_export = export_to_excel('Rzeczy i sprzęty', 0)
wb_to_export = export_to_excel('Hobby i przyjemności', 1)
wb_to_export = export_to_excel('Transport i noclegi', 2)
wb_to_export = export_to_excel('Podróże', 3)

wb_to_export.save(results_dir + '/' + part_label
                  + ' - zestawienie wydatków.xlsx')