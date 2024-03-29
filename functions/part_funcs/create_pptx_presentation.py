from pptx.util import Inches
from pptx import Presentation


def create_pptx_presentation(part_label, results_dir):

    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    blank_slide_layout = prs.slide_layouts[6]
    slides = list()

    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].shapes.title
    title.text = part_label + " - raport finansowy"

    # Part as a whole
    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].placeholders[1]
    title.text = "1. Cały przedział"

    for i in range(6):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(1.5)
        top = Inches(0.0)
        height = width = Inches(7.5)

        pic_path = results_dir + "/plots/plot" + str(i + 1) + ".png"
        slides[-1].shapes.add_picture(pic_path, left, top, height, width)

    # Averaged month
    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].placeholders[1]
    title.text = "2. Uśredniony miesiąc"

    for i in range(4):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(1.5)
        top = Inches(0.0)
        height = width = Inches(7.5)

        pic_path = results_dir + "/plots/plot" + str(i + 1 + 6) + ".png"
        slides[-1].shapes.add_picture(pic_path, left, top, height, width)

    # Part as a sequence of months
    slides.append(prs.slides.add_slide(title_slide_layout))
    title = slides[-1].placeholders[1]
    title.text = "3. Sekwencja miesięcy"

    for i in range(5):
        slides.append(prs.slides.add_slide(blank_slide_layout))
        left = Inches(0.0)
        top = Inches(0.1)
        height = Inches(10.5)
        width = Inches(7.0)

        pic_path = results_dir + "/plots/plot" + str(i + 1 + 10) + ".png"
        slides[-1].shapes.add_picture(pic_path, left, top, height, width)

    prs.save(results_dir + "/" + part_label + " - raport finansowy.pptx")
